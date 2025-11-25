from docx import Document
from pptx import Presentation
import io

def export_docx(project, sections):
    doc = Document()
    doc.add_heading(project.title, 0)
    
    # Sort sections by order
    sorted_sections = sorted(sections, key=lambda x: x.order_index)
    
    for section in sorted_sections:
        doc.add_heading(section.title, level=1)
        doc.add_paragraph(section.content)
        
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

def export_pptx(project, sections):
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = project.title
    slide.placeholders[1].text = project.topic
    
    # Content Slides
    bullet_slide_layout = prs.slide_layouts[1]
    sorted_sections = sorted(sections, key=lambda x: x.order_index)
    
    for section in sorted_sections:
        slide = prs.slides.add_slide(bullet_slide_layout)
        slide.shapes.title.text = section.title
        slide.placeholders[1].text = section.content
        
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer